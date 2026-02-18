"""
Document processing service for extracting text from various file formats.
Supports PDF, DOCX, PPTX, and TXT files.
"""

import logging
import re
import base64
from typing import Tuple, Optional, Dict, Any
from io import BytesIO

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
import chardet
import httpx

from app.models import DocumentType

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """
    Handles text extraction from various document formats.
    Optimized for memory efficiency on free-tier cloud platforms.
    """
    
    def __init__(self):
        self.supported_types = {
            DocumentType.PDF: self._extract_pdf,
            DocumentType.DOCX: self._extract_docx,
            DocumentType.PPTX: self._extract_pptx,
            DocumentType.TXT: self._extract_txt,
        }
    
    async def download_from_s3(self, s3_url: str) -> Tuple[bytes, Optional[str]]:
        """
        Download file content from S3 URL.
        
        Args:
            s3_url: The S3 URL to download from
            
        Returns:
            Tuple of (file_content, error_message)
        """
        try:
            async with httpx.AsyncClient(timeout=60.0) as client:
                logger.info(f"Downloading file from S3: {s3_url}")
                response = await client.get(s3_url)
                
                if response.status_code != 200:
                    error_msg = f"Failed to download from S3: HTTP {response.status_code}"
                    logger.error(error_msg)
                    return b"", error_msg
                
                content = response.content
                logger.info(f"Downloaded {len(content)} bytes from S3")
                return content, None
                
        except httpx.TimeoutException:
            error_msg = "Timeout while downloading from S3"
            logger.error(error_msg)
            return b"", error_msg
        except Exception as e:
            error_msg = f"Error downloading from S3: {str(e)}"
            logger.error(error_msg)
            return b"", error_msg
    
    async def extract_text(
        self, 
        file_content: bytes, 
        document_type: DocumentType
    ) -> Tuple[str, Optional[str]]:
        """
        Extract text from document.
        
        Args:
            file_content: Raw file bytes
            document_type: Type of document
            
        Returns:
            Tuple of (extracted_text, error_message)
        """
        try:
            extractor = self.supported_types.get(document_type)
            if not extractor:
                return "", f"Unsupported document type: {document_type}"
            
            text = extractor(file_content)
            cleaned_text = self._clean_text(text)
            
            if not cleaned_text.strip():
                return "", "No text content could be extracted from the document"
            
            logger.info(
                f"Extracted {len(cleaned_text)} characters from {document_type.value} document"
            )
            return cleaned_text, None
            
        except Exception as e:
            logger.error(f"Error extracting text from {document_type.value}: {str(e)}")
            return "", f"Failed to process document: {str(e)}"
    
    def _extract_pdf(self, content: bytes) -> str:
        """Extract text from PDF using PyMuPDF."""
        text_parts = []
        
        with fitz.open(stream=content, filetype="pdf") as doc:
            for page_num, page in enumerate(doc, 1):
                # Extract text with layout preservation
                page_text = page.get_text("text")
                
                if page_text.strip():
                    # Add page marker for section detection
                    text_parts.append(f"\n[Page {page_num}]\n{page_text}")
                
                # Also try to extract text from tables
                tables = page.find_tables()
                for table in tables:
                    table_text = self._format_table(table)
                    if table_text:
                        text_parts.append(table_text)
        
        return "\n".join(text_parts)
    
    def _format_table(self, table) -> str:
        """Format table data as readable text."""
        try:
            rows = table.extract()
            if not rows:
                return ""
            
            formatted_rows = []
            for row in rows:
                # Filter None values and join cells
                cells = [str(cell) if cell else "" for cell in row]
                formatted_rows.append(" | ".join(cells))
            
            return "\n".join(formatted_rows)
        except Exception:
            return ""
    
    def _extract_docx(self, content: bytes) -> str:
        """Extract text from DOCX preserving structure."""
        text_parts = []
        doc = DocxDocument(BytesIO(content))
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Check if it's a heading
                if para.style and para.style.name.startswith('Heading'):
                    level = para.style.name.replace('Heading ', '')
                    text_parts.append(f"\n{'#' * int(level) if level.isdigit() else '##'} {text}\n")
                else:
                    text_parts.append(text)
        
        # Also extract text from tables
        for table in doc.tables:
            table_text = self._extract_docx_table(table)
            if table_text:
                text_parts.append(table_text)
        
        return "\n".join(text_parts)
    
    def _extract_docx_table(self, table) -> str:
        """Extract text from DOCX table."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)
    
    def _extract_pptx(self, content: bytes) -> str:
        """Extract text from PowerPoint presentation."""
        text_parts = []
        prs = Presentation(BytesIO(content))
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            slide_text.append(f"\n[Slide {slide_num}]\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                
                # Handle tables in slides
                if shape.has_table:
                    table_text = self._extract_pptx_table(shape.table)
                    if table_text:
                        slide_text.append(table_text)
            
            if len(slide_text) > 1:  # More than just the slide marker
                text_parts.append("\n".join(slide_text))
        
        return "\n".join(text_parts)
    
    def _extract_pptx_table(self, table) -> str:
        """Extract text from PowerPoint table."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)
    
    def _extract_txt(self, content: bytes) -> str:
        """Extract text from plain text file with encoding detection."""
        # Detect encoding
        detected = chardet.detect(content)
        encoding = detected.get('encoding', 'utf-8') or 'utf-8'
        
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            # Fallback to utf-8 with error handling
            return content.decode('utf-8', errors='replace')
    
    def _clean_text(self, text: str) -> str:
        """
        Clean extracted text by removing noise and normalizing whitespace.
        """
        if not text:
            return ""
        
        # Remove common header/footer patterns
        patterns_to_remove = [
            r'Page \d+ of \d+',
            r'^\d+$',  # Standalone page numbers
            r'Confidential',
            r'All Rights Reserved',
            r'Copyright ©.*',
        ]
        
        for pattern in patterns_to_remove:
            text = re.sub(pattern, '', text, flags=re.MULTILINE | re.IGNORECASE)
        
        # Normalize whitespace
        text = re.sub(r'\n{3,}', '\n\n', text)  # Max 2 consecutive newlines
        text = re.sub(r'[ \t]+', ' ', text)  # Collapse spaces/tabs
        text = re.sub(r' +\n', '\n', text)  # Remove trailing spaces
        
        # Fix common OCR/extraction issues
        text = re.sub(r'(\w)-\n(\w)', r'\1\2', text)  # Fix hyphenated line breaks
        
        return text.strip()
    
    async def process_attachment(
        self,
        s3_url: str,
        attachment_type: str,
        file_size: int,
        text_threshold: int,
        file_size_threshold: int
    ) -> Dict[str, Any]:
        """
        Process attachment file with intelligent flow decision.
        
        Flow 1 (Direct Injection): Small files - extract text/base64 and return content
        Flow 2 (Vector Storage): Large files - return file content for chunking/embedding
        
        Args:
            s3_url: S3 URL of the attachment
            attachment_type: Type of attachment (IMAGE, PDF, DOCX, PPTX, TXT)
            file_size: File size in bytes
            text_threshold: Character count threshold for Flow 1
            file_size_threshold: File size threshold in bytes
            
        Returns:
            Dict with:
                - success: bool
                - flow: "direct_injection" or "vector_storage"
                - extracted_content: str (for Flow 1) or None
                - file_content: bytes (for Flow 2) or None
                - error: str or None
                - char_count: int (extracted text length)
        """
        try:
            # Download file from S3
            file_content, error = await self.download_from_s3(s3_url)
            if error:
                return {
                    "success": False,
                    "flow": None,
                    "extracted_content": None,
                    "file_content": None,
                    "error": error,
                    "char_count": 0
                }
            
            # Process based on attachment type
            if attachment_type.upper() == "IMAGE":
                # Convert image to base64
                extracted_content = base64.b64encode(file_content).decode('utf-8')
                # Prefix with data URI scheme for images
                extracted_content = f"data:image/*;base64,{extracted_content}"
                char_count = len(extracted_content)
                
            else:
                # Extract text from document
                doc_type_map = {
                    "PDF": DocumentType.PDF,
                    "DOCX": DocumentType.DOCX,
                    "PPTX": DocumentType.PPTX,
                    "TXT": DocumentType.TXT
                }
                
                doc_type = doc_type_map.get(attachment_type.upper())
                if not doc_type:
                    return {
                        "success": False,
                        "flow": None,
                        "extracted_content": None,
                        "file_content": None,
                        "error": f"Unsupported attachment type: {attachment_type}",
                        "char_count": 0
                    }
                
                extracted_content, extract_error = await self.extract_text(
                    file_content, 
                    doc_type
                )
                
                if extract_error:
                    return {
                        "success": False,
                        "flow": None,
                        "extracted_content": None,
                        "file_content": None,
                        "error": extract_error,
                        "char_count": 0
                    }
                
                char_count = len(extracted_content)
            
            # Decision logic: BOTH conditions must be met for Flow 1
            use_direct_injection = (
                char_count <= text_threshold and 
                file_size <= file_size_threshold
            )
            
            if use_direct_injection:
                logger.info(
                    f"Flow 1 (Direct Injection): {char_count} chars, "
                    f"{file_size} bytes - under thresholds"
                )
                return {
                    "success": True,
                    "flow": "direct_injection",
                    "extracted_content": extracted_content,
                    "file_content": None,
                    "error": None,
                    "char_count": char_count
                }
            else:
                logger.info(
                    f"Flow 2 (Vector Storage): {char_count} chars, "
                    f"{file_size} bytes - exceeds thresholds "
                    f"(text: {text_threshold}, size: {file_size_threshold})"
                )
                return {
                    "success": True,
                    "flow": "vector_storage",
                    "extracted_content": extracted_content,  # For chunking
                    "file_content": None,
                    "error": None,
                    "char_count": char_count
                }
                
        except Exception as e:
            logger.error(f"Error processing attachment: {str(e)}", exc_info=True)
            return {
                "success": False,
                "flow": None,
                "extracted_content": None,
                "file_content": None,
                "error": f"Failed to process attachment: {str(e)}",
                "char_count": 0
            }


# Singleton instance
document_processor = DocumentProcessor()
